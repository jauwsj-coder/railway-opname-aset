import json
import os
import base64
import csv
import io
import urllib.error
import urllib.request
from datetime import datetime, time
from zoneinfo import ZoneInfo

import gspread
from flask import Flask, Response, jsonify, render_template, request
from google.oauth2.service_account import Credentials
from itsdangerous import BadSignature, SignatureExpired, URLSafeTimedSerializer
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


app = Flask(__name__)
MASTER_SHEET, LOG_SHEET, ROLE_SHEET, DASHBOARD_SHEET = "MASTER_ASET", "LOG_OPNAME", "ROLE", "DASHBOARD"
TIMEZONE = os.getenv("APP_TIMEZONE", "Asia/Jakarta")
ALLOWED_IMAGE_TYPES = {"image/jpeg", "image/png", "image/webp"}
AUTH_MAX_AGE = 12 * 60 * 60

VALID_ROLES = {"SUPER ADMIN", "SUPER ADMIN, PIC ASET", "PIC ASET"}
PIC_ROLES = {"SUPER ADMIN, PIC ASET", "PIC ASET"}
GOOD_CONDITIONS = {"OK", "BAIK", "GOOD"}
DAMAGED_CONDITIONS = {"RUSAK", "BROKEN", "MAINTENANCE", "NOT OK"}

MASTER_HEADERS = ["NOMOR ASSET", "TYPE", "NO LAYOUT", "USER", "OPNAME", "KONDISI", "LOKASI DETAIL", "AREA", "KONDISI TERAKHIR", "STATUS TERAKHIR", "TANGGAL OPNAME TERAKHIR", "KETERANGAN TERAKHIR", "DOKUMENTASI TERAKHIR"]
LOG_HEADERS = ["TIMESTAMP", "NOMOR ASSET", "TYPE", "NO LAYOUT", "USER", "OPNAME", "KONDISI", "LOKASI DETAIL", "AREA", "KONDISI TERAKHIR", "STATUS TERAKHIR", "TANGGAL OPNAME TERAKHIR", "KETERANGAN TERAKHIR", "DOKUMENTASI TERAKHIR", "NAMA PETUGAS", "ID USER", "ROLE"]
ROLE_HEADERS = ["NAMA USER", "ID USER", "ROLE", "AREA", "AREA SCORECARD"]
DASHBOARD_HEADERS = ["METRIK", "NILAI", "DIPERBARUI"]
SCORE_HEADERS = ["NAMA PETUGAS", "ID USER", "ROLE", "AREA SCORECARD", "TOTAL ASSET", "SUDAH OPNAME", "BELUM OPNAME", "PROGRESS", "STATUS", "DOKUMENTASI ADA", "KETERANGAN ADA", "ASET BAIK", "ASET RUSAK"]
DATA_QUALITY_CATEGORIES = {
    "duplicate_asset": "Nomor Aset Double",
    "missing_asset_number": "Belum Ada Nomor Aset",
    "empty_area": "Area Kosong",
    "empty_location": "Lokasi Detail Kosong",
    "empty_type": "Type Kosong",
    "empty_user": "User/PIC Kosong",
    "not_opnamed": "Belum Opname",
    "empty_documentation": "Dokumentasi Kosong",
    "damaged_without_notes": "Aset Rusak Tanpa Keterangan",
    "completed_opname": "Sudah Opname",
    "damaged_with_notes": "Aset Rusak dengan Keterangan",
}
DATA_QUALITY_EXPORT_FIELDS = ["SUMBER", "BARIS", "NOMOR ASSET", "TYPE", "USER", "AREA", "LOKASI DETAIL", "KONDISI", "STATUS", "KETERANGAN", "MASALAH"]
MISSING_ASSET_MARKERS = {
    "TIDAK ADA NOMOR ASET", "TIDAK ADA NOMOR ASSET", "BELUM ADA NOMOR ASET",
    "BELUM ADA NOMOR ASSET", "NO ASSET KOSONG",
}


class AppError(Exception):
    def __init__(self, message, status=400):
        super().__init__(message)
        self.message, self.status = message, status


@app.errorhandler(AppError)
def handle_app_error(error):
    return jsonify({"success": False, "message": error.message}), error.status


@app.errorhandler(Exception)
def handle_unexpected_error(error):
    app.logger.exception(error)
    return jsonify({"success": False, "message": "Terjadi kesalahan pada server."}), 500


@app.errorhandler(404)
def handle_not_found(error):
    return jsonify({"success": False, "message": "Endpoint tidak ditemukan.", "path": request.path}), 404


@app.get("/")
def index():
    return render_template("index.html")


@app.get("/healthz")
def healthz():
    return jsonify({"status": "ok"})


@app.get("/api/users")
def users():
    rows = get_rows(get_worksheet(ROLE_SHEET), ROLE_HEADERS)
    names = sorted({clean(row["NAMA USER"]) for row in rows if clean(row["NAMA USER"])})
    return jsonify({"users": names})


@app.post("/api/login")
def login():
    payload = request.get_json(silent=True) or {}
    user = find_role_user(payload.get("name"), payload.get("userId"))
    if not user:
        raise AppError("NAMA USER dan ID USER tidak cocok.", 401)
    identity = validated_identity(user)
    ensure_pic_has_assets(identity)
    return jsonify({"token": serializer().dumps(identity), "user": identity})


@app.get("/api/dashboard")
def dashboard():
    identity = require_user()
    start_date, end_date = parse_period(request.args.get("startDate"), request.args.get("endDate"))
    score_start, score_end, score_period = parse_score_period(request.args.get("scorePeriod"))
    summary, score, warnings = build_dashboard(identity, start_date, end_date, score_start, score_end)
    return jsonify({"summary": summary, "scoreCard": score, "warnings": warnings, "scope": identity["area"], "scorePeriod": score_period})


@app.post("/api/dashboard/sync")
def sync_dashboard():
    identity = require_user()
    if identity["role"] not in {"SUPER ADMIN", "SUPER ADMIN, PIC ASET"} or not can_access_all(identity):
        raise AppError("Hanya SUPER ADMIN dengan AREA ALL yang dapat melakukan Sync Sheet.", 403)
    score_start, score_end, _ = parse_score_period(request.args.get("scorePeriod"))
    summary, score, _ = build_dashboard(all_area_identity(), score_start_date=score_start, score_end_date=score_end)
    write_dashboard_sheet(summary, score)
    return jsonify({"success": True, "message": "Sheet DASHBOARD berhasil diperbarui."})


@app.get("/api/data-quality")
def data_quality_summary():
    identity = require_data_quality_access()
    start_date, end_date, period = parse_score_period(request.args.get("period"))
    results, warnings = build_data_quality(identity, start_date, end_date)
    return jsonify({
        "summary": [{"key": key, "label": DATA_QUALITY_CATEGORIES[key], "count": len(results[key])} for key in DATA_QUALITY_CATEGORIES],
        "warnings": warnings,
        "period": period,
    })


@app.get("/api/data-quality/detail")
def data_quality_detail():
    identity = require_data_quality_access()
    category = clean(request.args.get("category"))
    if category not in DATA_QUALITY_CATEGORIES:
        raise AppError("Kategori pemeriksaan data tidak valid.")
    start_date, end_date, period = parse_score_period(request.args.get("period"))
    results, warnings = build_data_quality(identity, start_date, end_date)
    return jsonify({"category": category, "label": DATA_QUALITY_CATEGORIES[category], "rows": results[category], "warnings": warnings, "period": period})


@app.get("/api/data-quality/export")
def data_quality_export():
    identity = require_data_quality_access()
    category = clean(request.args.get("category"))
    if category not in DATA_QUALITY_CATEGORIES:
        raise AppError("Kategori pemeriksaan data tidak valid.")
    start_date, end_date, period = parse_score_period(request.args.get("period"))
    results, _ = build_data_quality(identity, start_date, end_date)
    export_format = normalize(request.args.get("format")) or "CSV"
    rows, label = results[category], DATA_QUALITY_CATEGORIES[category]
    filename_base = f"pemeriksaan-data-{category}-{period.lower()}"
    if export_format == "CSV":
        output = io.StringIO()
        writer = csv.DictWriter(output, fieldnames=DATA_QUALITY_EXPORT_FIELDS)
        writer.writeheader()
        writer.writerows(rows)
        return download_response(("\ufeff" + output.getvalue()).encode("utf-8"), "text/csv; charset=utf-8", filename_base + ".csv")
    if export_format == "PDF":
        return download_response(build_data_quality_pdf(label, period, rows), "application/pdf", filename_base + ".pdf")
    if export_format in {"PPT", "PPTX"}:
        return download_response(build_data_quality_ppt(label, period, rows), "application/vnd.openxmlformats-officedocument.presentationml.presentation", filename_base + ".pptx")
    raise AppError("Format export harus CSV, PDF, atau PPTX.")


@app.get("/api/assets/<asset_code>")
def find_asset(asset_code):
    identity = require_user()
    code = normalize(asset_code)
    if not code:
        raise AppError("NOMOR ASSET wajib diisi.")
    asset = next((row for row in get_rows(get_worksheet(MASTER_SHEET), MASTER_HEADERS) if normalize(row["NOMOR ASSET"]) == code), None)
    if not asset:
        raise AppError(f"Aset {code} tidak ditemukan.", 404)
    ensure_area_access(identity, asset["AREA"])
    history, warnings = [], []
    try:
        history = [row for row in get_rows(get_worksheet(LOG_SHEET), LOG_HEADERS) if normalize(row["NOMOR ASSET"]) == code and can_access_area(identity, row["AREA"])]
        history.reverse()
    except AppError as error:
        warnings.append(f"Riwayat opname belum dapat ditampilkan: {error.message}")
    return jsonify({"asset": serialize_asset(asset), "history": [serialize_log(row) for row in history[:5]], "warnings": warnings})


@app.post("/api/opname")
def submit_opname():
    operator = require_user()
    payload = request.get_json(silent=True) or {}
    code, condition = normalize(payload.get("assetCode")), normalize(payload.get("condition"))
    notes, documentation = clean(payload.get("notes")), clean(payload.get("documentation"))
    if not code:
        raise AppError("NOMOR ASSET wajib diisi.")
    if condition not in GOOD_CONDITIONS | DAMAGED_CONDITIONS:
        raise AppError("KONDISI TERAKHIR harus OK/BAIK/GOOD atau RUSAK/BROKEN/MAINTENANCE/NOT OK.")

    master = get_worksheet(MASTER_SHEET)
    values = get_sheet_values(master, MASTER_SHEET)
    validate_headers(values, MASTER_HEADERS, MASTER_SHEET)
    headers = values[0]
    header_map = {header: index + 1 for index, header in enumerate(headers)}
    asset_row, asset = None, None
    for row_number, row_values in enumerate(values[1:], start=2):
        row = row_to_dict(headers, row_values)
        if normalize(row["NOMOR ASSET"]) == code:
            asset_row, asset = row_number, row
            break
    if not asset:
        raise AppError(f"Aset {code} tidak ditemukan.", 404)
    ensure_area_access(operator, asset["AREA"])

    date_value, status, opname_value = now_text(), current_period_status(), "DONE"
    log = get_worksheet(LOG_SHEET)
    ensure_required_headers(log, LOG_SHEET, LOG_HEADERS)
    append_record(log, LOG_SHEET, LOG_HEADERS, {
        "TIMESTAMP": date_value, "NOMOR ASSET": asset["NOMOR ASSET"], "TYPE": asset["TYPE"],
        "NO LAYOUT": asset["NO LAYOUT"], "USER": asset["USER"], "OPNAME": opname_value,
        "KONDISI": condition, "LOKASI DETAIL": asset["LOKASI DETAIL"], "AREA": asset["AREA"],
        "KONDISI TERAKHIR": condition, "STATUS TERAKHIR": status, "TANGGAL OPNAME TERAKHIR": date_value,
        "KETERANGAN TERAKHIR": notes, "NAMA PETUGAS": operator["name"], "ID USER": operator["userId"],
        "ROLE": operator["role"], "DOKUMENTASI TERAKHIR": documentation,
    })

    master.batch_update([
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["OPNAME"]), "values": [[opname_value]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["KONDISI"]), "values": [[condition]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["KONDISI TERAKHIR"]), "values": [[condition]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["STATUS TERAKHIR"]), "values": [[status]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["TANGGAL OPNAME TERAKHIR"]), "values": [[date_value]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["KETERANGAN TERAKHIR"]), "values": [[notes]]},
        {"range": gspread.utils.rowcol_to_a1(asset_row, header_map["DOKUMENTASI TERAKHIR"]), "values": [[documentation]]},
    ], value_input_option="USER_ENTERED")

    summary, score, warnings = build_dashboard(operator)
    global_summary, global_score, _ = build_dashboard(all_area_identity())
    try:
        write_dashboard_sheet(global_summary, global_score)
    except AppError as error:
        warnings.append(f"Opname tersimpan, tetapi sheet DASHBOARD belum diperbarui: {error.message}")
    return jsonify({"success": True, "message": "Opname aset berhasil disimpan.", "summary": summary, "scoreCard": score, "warnings": warnings})


@app.post("/api/upload-documentation")
def upload_documentation():
    identity = require_user()
    file = request.files.get("photo")
    asset_code = normalize(request.form.get("assetCode"))
    if not file or not file.filename:
        return jsonify({"success": True, "url": "", "message": "Tidak ada foto yang diunggah."})
    if file.mimetype not in ALLOWED_IMAGE_TYPES:
        raise AppError("Format foto harus JPG, PNG, atau WEBP.")
    if not asset_code:
        raise AppError("NOMOR ASSET wajib diisi sebelum upload.")
    file.stream.seek(0, os.SEEK_END)
    size = file.stream.tell()
    file.stream.seek(0)
    if size > 10 * 1024 * 1024:
        raise AppError("Ukuran foto maksimal 10 MB.")
    extension = os.path.splitext(file.filename)[1].lower() or ".jpg"
    name = f"{asset_code}_{datetime.now(ZoneInfo(TIMEZONE)).strftime('%Y%m%d_%H%M%S')}_{normalize(identity['userId'])}{extension}"
    result = call_photo_upload_script({
        "action": "upload",
        "fileName": name,
        "mimeType": file.mimetype,
        "base64Data": base64.b64encode(file.read()).decode("ascii"),
        "assetCode": asset_code,
    })
    return jsonify({"success": True, **result})


@app.get("/api/cleanup-drive-photos")
def cleanup_drive_photos_endpoint():
    expected_token = os.getenv("SETUP_TOKEN", "")
    if not expected_token or request.args.get("token", "") != expected_token:
        raise AppError("Setup token tidak valid.", 403)
    return jsonify({"success": True, **call_photo_upload_script({"action": "cleanup"})})


@app.get("/api/test-photo-upload")
def test_photo_upload_endpoint():
    expected_token = os.getenv("SETUP_TOKEN", "")
    if not expected_token or request.args.get("token", "") != expected_token:
        raise AppError("Setup token tidak valid.", 403)
    probe_png = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
    result = call_photo_upload_script({
        "action": "test",
        "fileName": f"railway-upload-test-{datetime.now().strftime('%Y%m%d%H%M%S')}.png",
        "mimeType": "image/png",
        "base64Data": probe_png,
    })
    return jsonify({"success": True, **result})


@app.post("/api/setup")
def setup_sheets():
    if request.headers.get("X-Setup-Token", "") != os.getenv("SETUP_TOKEN", "") or not os.getenv("SETUP_TOKEN"):
        raise AppError("Setup token tidak valid.", 403)
    spreadsheet = get_spreadsheet()
    results = [ensure_worksheet(spreadsheet, name, headers) for name, headers in (
        (MASTER_SHEET, MASTER_HEADERS), (LOG_SHEET, LOG_HEADERS), (ROLE_SHEET, ROLE_HEADERS), (DASHBOARD_SHEET, DASHBOARD_HEADERS)
    )]
    return jsonify({"success": True, "message": "Setup selesai tanpa menghapus data existing.", "sheets": results})


def serializer():
    secret = os.getenv("APP_SECRET_KEY", "").strip()
    if not secret:
        raise AppError("APP_SECRET_KEY belum dikonfigurasi.", 503)
    return URLSafeTimedSerializer(secret, salt="opname-user")


def require_user():
    header = request.headers.get("Authorization", "")
    if not header.startswith("Bearer "):
        raise AppError("Silakan masuk terlebih dahulu.", 401)
    try:
        token_identity = serializer().loads(header[7:], max_age=AUTH_MAX_AGE)
    except (BadSignature, SignatureExpired) as exc:
        raise AppError("Sesi sudah tidak valid. Silakan masuk kembali.", 401) from exc
    user = find_role_user(token_identity.get("name"), token_identity.get("userId"))
    if not user:
        raise AppError("User tidak lagi terdaftar pada sheet ROLE.", 401)
    identity = validated_identity(user)
    ensure_pic_has_assets(identity)
    return identity


def find_role_user(name, user_id):
    name_key, id_key = normalize(name), normalize(user_id)
    return next((row for row in get_rows(get_worksheet(ROLE_SHEET), ROLE_HEADERS) if normalize(row["NAMA USER"]) == name_key and normalize(row["ID USER"]) == id_key), None)


def validated_identity(row):
    role, area = normalize(row["ROLE"]), normalize(row["AREA"])
    if role not in VALID_ROLES:
        raise AppError(f"ROLE tidak valid: {clean(row['ROLE']) or '-'}", 403)
    if not area:
        raise AppError(f"AREA kosong untuk user {clean(row['NAMA USER']) or '-'}.", 403)
    return {"name": clean(row["NAMA USER"]), "userId": clean(row["ID USER"]), "role": role, "area": area, "areas": sorted(parse_areas(area))}


def ensure_pic_has_assets(identity):
    if identity["role"] != "PIC ASET":
        return
    assets = get_rows(get_worksheet(MASTER_SHEET), MASTER_HEADERS)
    if not any(normalize(row["NOMOR ASSET"]) and can_access_area(identity, row["AREA"]) for row in assets):
        raise AppError(f"PIC ASET tidak memiliki aset sesuai AREA {identity['area']}.", 403)


def require_data_quality_access():
    identity = require_user()
    if identity["role"] not in {"SUPER ADMIN", "SUPER ADMIN, PIC ASET"}:
        raise AppError("Menu Pemeriksaan Data hanya dapat diakses SUPER ADMIN.", 403)
    return all_area_identity()


def build_data_quality(identity, start_date=None, end_date=None):
    warnings = []
    master_rows = get_rows_with_numbers(get_worksheet(MASTER_SHEET), MASTER_HEADERS)
    log_available = True
    try:
        log_rows = get_rows_with_numbers(get_worksheet(LOG_SHEET), LOG_HEADERS)
    except AppError as error:
        log_available = False
        log_rows = []
        warnings.append(f"Pemeriksaan berbasis LOG_OPNAME belum dapat dihitung: {error.message}")

    master_rows = [row for row in master_rows if can_access_area(identity, row["AREA"]) or not clean(row["AREA"])]
    log_rows = [row for row in log_rows if can_access_area(identity, row["AREA"]) and log_in_period(row, start_date, end_date)]
    results = {key: [] for key in DATA_QUALITY_CATEGORIES}

    code_groups = {}
    for row in master_rows:
        code = normalize(row["NOMOR ASSET"])
        if code and code not in MISSING_ASSET_MARKERS:
            code_groups.setdefault(code, []).append(row)
    for rows in code_groups.values():
        if len(rows) > 1:
            for row in rows:
                results["duplicate_asset"].append(data_quality_record(row, "MASTER_ASET", f"NOMOR ASSET muncul {len(rows)} kali di MASTER_ASET."))

    valid_master_codes = set(code_groups)
    log_codes = {normalize(row["NOMOR ASSET"]) for row in log_rows if normalize(row["NOMOR ASSET"])}
    for row in master_rows:
        code = normalize(row["NOMOR ASSET"])
        if not code or code in MISSING_ASSET_MARKERS:
            results["missing_asset_number"].append(data_quality_record(row, "MASTER_ASET", "NOMOR ASSET kosong atau menggunakan penanda belum ada nomor aset."))
        if not clean(row["AREA"]):
            results["empty_area"].append(data_quality_record(row, "MASTER_ASET", "AREA kosong."))
        if not clean(row["LOKASI DETAIL"]):
            results["empty_location"].append(data_quality_record(row, "MASTER_ASET", "LOKASI DETAIL kosong."))
        if not clean(row["TYPE"]):
            results["empty_type"].append(data_quality_record(row, "MASTER_ASET", "TYPE kosong."))
        if not clean(row["USER"]):
            results["empty_user"].append(data_quality_record(row, "MASTER_ASET", "USER/PIC kosong."))
        if log_available and code in valid_master_codes and code not in log_codes:
            results["not_opnamed"].append(data_quality_record(row, "MASTER_ASET", "NOMOR ASSET belum muncul di LOG_OPNAME pada periode terpilih."))

    for row in log_rows:
        if is_completed_log(row) and not clean(row["DOKUMENTASI TERAKHIR"]):
            results["empty_documentation"].append(data_quality_record(row, "LOG_OPNAME", "Data opname belum memiliki dokumentasi."))
        if log_condition(row) in DAMAGED_CONDITIONS and not clean(row["KETERANGAN TERAKHIR"]):
            results["damaged_without_notes"].append(data_quality_record(row, "LOG_OPNAME", "Aset rusak belum memiliki keterangan."))
    latest_logs = {}
    for row in log_rows:
        code = normalize(row["NOMOR ASSET"])
        if code:
            latest_logs[code] = row
    for row in latest_logs.values():
        if is_completed_log(row):
            results["completed_opname"].append(data_quality_record(row, "LOG_OPNAME", "Aset sudah diopname pada periode terpilih."))
        if log_condition(row) in DAMAGED_CONDITIONS and clean(row["KETERANGAN TERAKHIR"]):
            results["damaged_with_notes"].append(data_quality_record(row, "LOG_OPNAME", "Aset rusak memiliki keterangan dan perlu ditindaklanjuti."))
    return results, warnings


def data_quality_record(row, source, problem):
    return {
        "NOMOR ASSET": clean(row["NOMOR ASSET"]) or "-",
        "TYPE": clean(row["TYPE"]) or "-",
        "USER": clean(row["USER"]) or "-",
        "AREA": clean(row["AREA"]) or "-",
        "LOKASI DETAIL": clean(row["LOKASI DETAIL"]) or "-",
        "KONDISI": clean(row.get("KONDISI TERAKHIR") or row.get("KONDISI")) or "-",
        "STATUS": clean(row.get("STATUS TERAKHIR") or row.get("OPNAME")) or "-",
        "KETERANGAN": clean(row.get("KETERANGAN TERAKHIR")) or "-",
        "MASALAH": problem,
        "SUMBER": source,
        "BARIS": row.get("_ROW_NUMBER", ""),
    }


def is_completed_log(row):
    return normalize(row.get("OPNAME")) == "DONE" or "OPNAME" in normalize(row.get("STATUS TERAKHIR"))


def download_response(data, mimetype, filename):
    return Response(data, mimetype=mimetype, headers={"Content-Disposition": f'attachment; filename="{filename}"'})


def build_data_quality_pdf(label, period, rows):
    output = io.BytesIO()
    document = SimpleDocTemplate(output, pagesize=landscape(A4), rightMargin=10 * mm, leftMargin=10 * mm, topMargin=10 * mm, bottomMargin=10 * mm)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Laporan Pemeriksaan Data Aset", styles["Title"]),
        Paragraph(f"{label} | Periode: {period} | Jumlah: {len(rows)}", styles["Heading2"]),
        Spacer(1, 5 * mm),
    ]
    headers = ["NO ASSET", "TYPE", "USER", "AREA", "LOKASI", "KONDISI", "STATUS", "KETERANGAN", "MASALAH"]
    table_rows = [headers]
    for row in rows:
        table_rows.append([clean(row[field]) for field in ["NOMOR ASSET", "TYPE", "USER", "AREA", "LOKASI DETAIL", "KONDISI", "STATUS", "KETERANGAN", "MASALAH"]])
    if not rows:
        table_rows.append(["Tidak ada data", "", "", "", "", "", "", "", ""])
    table = Table(table_rows, repeatRows=1, colWidths=[30 * mm, 22 * mm, 25 * mm, 22 * mm, 30 * mm, 20 * mm, 28 * mm, 35 * mm, 45 * mm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#009FB2")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 6.5),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#DCE3E8")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F6F9FA")]),
    ]))
    story.append(table)
    document.build(story)
    return output.getvalue()


def build_data_quality_ppt(label, period, rows):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    add_ppt_title(title_slide, "Laporan Pemeriksaan Data Aset", f"{label} | Periode: {period} | Jumlah: {len(rows)}")
    headers = ["NO ASSET", "TYPE", "USER", "AREA", "KONDISI", "KETERANGAN", "MASALAH"]
    chunks = [rows[index:index + 7] for index in range(0, len(rows), 7)] or [[]]
    for page, chunk in enumerate(chunks, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        add_ppt_title(slide, label, f"Periode {period} | Halaman {page}/{len(chunks)}")
        table_shape = slide.shapes.add_table(len(chunk) + 1, len(headers), Inches(.35), Inches(1.35), Inches(12.63), Inches(5.65))
        table = table_shape.table
        for column, header in enumerate(headers):
            table.cell(0, column).text = header
        for row_index, row in enumerate(chunk, start=1):
            for column, field in enumerate(["NOMOR ASSET", "TYPE", "USER", "AREA", "KONDISI", "KETERANGAN", "MASALAH"]):
                table.cell(row_index, column).text = clean(row[field])
        style_ppt_table(table)
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def add_ppt_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(.45), Inches(.25), Inches(12.4), Inches(.55))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = ppt_rgb("007F91")
    subtitle_box = slide.shapes.add_textbox(Inches(.45), Inches(.82), Inches(12.4), Inches(.35))
    subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
    subtitle_paragraph.text = subtitle
    subtitle_paragraph.font.size = Pt(11)
    subtitle_paragraph.font.color.rgb = ppt_rgb("647180")


def style_ppt_table(table):
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(.04)
            cell.margin_right = Inches(.04)
            cell.margin_top = Inches(.03)
            cell.margin_bottom = Inches(.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ppt_rgb("009FB2" if row_index == 0 else "F6F9FA")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = ppt_rgb("FFFFFF" if row_index == 0 else "18222D")


def ppt_rgb(value):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(value)


def build_dashboard(identity, start_date=None, end_date=None, score_start_date=None, score_end_date=None):
    master_rows = get_rows(get_worksheet(MASTER_SHEET), MASTER_HEADERS)
    blank_area_codes = [clean(row["NOMOR ASSET"]) for row in master_rows if normalize(row["NOMOR ASSET"]) and not normalize(row["AREA"])]
    if blank_area_codes:
        raise AppError(f"AREA kosong pada MASTER_ASET untuk NOMOR ASSET: {', '.join(blank_area_codes[:5])}.", 503)
    assets = [row for row in master_rows if normalize(row["NOMOR ASSET"]) and can_access_area(identity, row["AREA"])]
    warnings = []
    try:
        logs = [row for row in get_rows(get_worksheet(LOG_SHEET), LOG_HEADERS) if normalize(row["NOMOR ASSET"]) and can_access_area(identity, row["AREA"]) and log_in_period(row, start_date, end_date)]
    except AppError as error:
        logs = []
        warnings.append(f"Data opname belum dapat dihitung: {error.message}")
    asset_codes = {normalize(row["NOMOR ASSET"]) for row in assets}
    scoped_logs = [row for row in logs if normalize(row["NOMOR ASSET"]) in asset_codes]
    latest_by_asset = {}
    for row in scoped_logs:
        latest_by_asset[normalize(row["NOMOR ASSET"])] = row
    completed = len(latest_by_asset)
    good = sum(1 for row in latest_by_asset.values() if log_condition(row) in GOOD_CONDITIONS)
    damaged = sum(1 for row in latest_by_asset.values() if log_condition(row) in DAMAGED_CONDITIONS)
    summary = {"total": len(assets), "completed": completed, "pending": max(len(assets) - completed, 0), "good": good, "damaged": damaged}

    score = build_score_card(identity, master_rows, score_start_date, score_end_date, warnings)
    return summary, score, warnings


def build_score_card(viewer, master_rows, start_date=None, end_date=None, warnings=None):
    warnings = warnings if warnings is not None else []
    try:
        role_rows = get_rows(get_worksheet(ROLE_SHEET), ROLE_HEADERS)
    except AppError as error:
        warnings.append(f"Score Board belum dapat dihitung: {error.message}")
        return []
    try:
        all_logs = [row for row in get_rows(get_worksheet(LOG_SHEET), LOG_HEADERS) if normalize(row["NOMOR ASSET"]) and log_in_period(row, start_date, end_date)]
    except AppError as error:
        all_logs = []
        if not any(error.message in warning for warning in warnings):
            warnings.append(f"Score Board belum dapat dihitung dari LOG_OPNAME: {error.message}")

    score = []
    for role_row in role_rows:
        role = normalize(role_row["ROLE"])
        if role not in PIC_ROLES:
            continue
        score_area_value = clean(role_row["AREA SCORECARD"]) or clean(role_row["AREA"])
        score_areas = parse_areas(score_area_value)
        if not score_areas:
            continue
        if not can_view_score_areas(viewer, score_areas):
            continue

        asset_codes = {normalize(row["NOMOR ASSET"]) for row in master_rows if normalize(row["NOMOR ASSET"]) and area_in_scope(row["AREA"], score_areas)}
        matching_logs = [row for row in all_logs if normalize(row["NOMOR ASSET"]) in asset_codes and area_in_scope(row["AREA"], score_areas)]
        latest_by_asset = {}
        for row in matching_logs:
            latest_by_asset[normalize(row["NOMOR ASSET"])] = row

        total, completed = len(asset_codes), len(latest_by_asset)
        progress = round(completed * 100 / total, 2) if total else 0
        score.append({
            "name": clean(role_row["NAMA USER"]),
            "userId": clean(role_row["ID USER"]),
            "role": role,
            "scoreAreas": score_area_value,
            "total": total,
            "completed": completed,
            "pending": max(total - completed, 0),
            "progress": progress,
            "status": progress_status(progress),
            "documentationCount": sum(1 for row in latest_by_asset.values() if clean(row["DOKUMENTASI TERAKHIR"])),
            "notesCount": sum(1 for row in latest_by_asset.values() if clean(row["KETERANGAN TERAKHIR"])),
            "good": sum(1 for row in latest_by_asset.values() if log_condition(row) in GOOD_CONDITIONS),
            "damaged": sum(1 for row in latest_by_asset.values() if log_condition(row) in DAMAGED_CONDITIONS),
        })
    return sorted(score, key=lambda item: (-item["progress"], item["name"]))


def parse_period(start_value, end_value):
    try:
        start_date = datetime.combine(datetime.strptime(start_value, "%Y-%m-%d").date(), time.min, ZoneInfo(TIMEZONE)) if start_value else None
        end_date = datetime.combine(datetime.strptime(end_value, "%Y-%m-%d").date(), time.max, ZoneInfo(TIMEZONE)) if end_value else None
    except ValueError as exc:
        raise AppError("Format periode tanggal tidak valid. Gunakan YYYY-MM-DD.") from exc
    if start_date and end_date and start_date > end_date:
        raise AppError("Tanggal awal tidak boleh lebih besar dari tanggal akhir.")
    return start_date, end_date


def parse_score_period(value):
    period = normalize(value) or "ALL"
    if period == "ALL":
        return None, None, "ALL"
    year = datetime.now(ZoneInfo(TIMEZONE)).year
    if period == "JAN-JUN":
        return datetime(year, 1, 1, tzinfo=ZoneInfo(TIMEZONE)), datetime.combine(datetime(year, 6, 30).date(), time.max, ZoneInfo(TIMEZONE)), "JAN-JUN"
    if period == "JUL-DES":
        return datetime(year, 7, 1, tzinfo=ZoneInfo(TIMEZONE)), datetime.combine(datetime(year, 12, 31).date(), time.max, ZoneInfo(TIMEZONE)), "JUL-DES"
    raise AppError("Periode Score Board tidak valid.")


def log_in_period(row, start_date, end_date):
    if not start_date and not end_date:
        return True
    value = clean(row["TANGGAL OPNAME TERAKHIR"] or row["TIMESTAMP"])
    try:
        log_date = datetime.strptime(value, "%Y-%m-%d %H:%M:%S").replace(tzinfo=ZoneInfo(TIMEZONE))
    except ValueError:
        return False
    return (not start_date or log_date >= start_date) and (not end_date or log_date <= end_date)


def log_condition(row):
    return normalize(row["KONDISI TERAKHIR"] or row["KONDISI"])


def can_access_all(identity):
    return "ALL" in identity.get("areas", parse_areas(identity.get("area")))


def all_area_identity():
    return {"role": "SUPER ADMIN", "area": "ALL", "areas": ["ALL"]}


def can_access_area(identity, asset_area):
    return can_access_all(identity) or normalize(asset_area) in identity.get("areas", parse_areas(identity.get("area")))


def parse_areas(value):
    return {normalize(area) for area in clean(value).split(",") if normalize(area)}


def area_in_scope(asset_area, score_areas):
    return "ALL" in score_areas or normalize(asset_area) in score_areas


def can_view_score_areas(viewer, score_areas):
    viewer_areas = set(viewer.get("areas", parse_areas(viewer.get("area"))))
    return can_access_all(viewer) or "ALL" not in score_areas and bool(viewer_areas & score_areas)


def progress_status(progress):
    if progress >= 100:
        return "Selesai"
    if progress >= 90:
        return "Hampir Selesai"
    if progress >= 75:
        return "On Track"
    if progress >= 50:
        return "Dalam Proses"
    if progress > 0:
        return "Perlu Dukungan"
    return "Belum Mulai"


def ensure_area_access(identity, asset_area):
    if not clean(asset_area):
        raise AppError("AREA aset kosong pada MASTER_ASET.", 403)
    if not can_access_area(identity, asset_area):
        raise AppError(f"Anda tidak memiliki akses untuk memproses aset AREA {clean(asset_area)}.", 403)


def write_dashboard_sheet(summary, score):
    sheet = get_worksheet(DASHBOARD_SHEET)
    updated = now_text()
    values = [DASHBOARD_HEADERS, ["TOTAL ASSET", summary["total"], updated], ["SUDAH OPNAME", summary["completed"], updated], ["BELUM OPNAME", summary["pending"], updated], ["ASET BAIK", summary["good"], updated], ["ASET RUSAK", summary["damaged"], updated], [], SCORE_HEADERS]
    values.extend([[
        item["name"], item["userId"], item["role"], item["scoreAreas"], item["total"], item["completed"],
        item["pending"], item["progress"] / 100, item["status"], item["documentationCount"], item["notesCount"],
        item["good"], item["damaged"]
    ] for item in score])
    sheet.clear()
    sheet.update(values, "A1", value_input_option="USER_ENTERED")
    if score:
        sheet.format(f"H9:H{8 + len(score)}", {"numberFormat": {"type": "PERCENT", "pattern": "0.00%"}})


def get_spreadsheet():
    spreadsheet_id, credentials_json = os.getenv("GOOGLE_SHEET_ID", "").strip(), os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON", "").strip()
    if not spreadsheet_id or not credentials_json:
        raise AppError("Konfigurasi Google Sheets belum lengkap.", 503)
    try:
        return gspread.authorize(get_google_credentials()).open_by_key(spreadsheet_id)
    except json.JSONDecodeError as exc:
        raise AppError("GOOGLE_SERVICE_ACCOUNT_JSON bukan JSON yang valid.", 503) from exc
    except gspread.SpreadsheetNotFound as exc:
        raise AppError("Google Sheet tidak ditemukan atau service account belum memiliki akses.", 503) from exc
    except gspread.exceptions.APIError as exc:
        raise AppError("Google Sheet tidak bisa dibaca. Periksa API dan akses service account.", 503) from exc
    except (ValueError, KeyError) as exc:
        raise AppError("Kredensial service account tidak valid.", 503) from exc


def get_google_credentials():
    credentials_json = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON", "").strip()
    if not credentials_json:
        raise AppError("GOOGLE_SERVICE_ACCOUNT_JSON belum dikonfigurasi.", 503)
    try:
        info = json.loads(credentials_json)
        return Credentials.from_service_account_info(info, scopes=["https://www.googleapis.com/auth/spreadsheets"])
    except (json.JSONDecodeError, ValueError, KeyError) as exc:
        raise AppError("Kredensial service account tidak valid.", 503) from exc


def call_photo_upload_script(payload):
    script_url = os.getenv("PHOTO_UPLOAD_SCRIPT_URL", "").strip()
    secret = os.getenv("PHOTO_UPLOAD_SECRET", "").strip()
    if not script_url:
        raise AppError("PHOTO_UPLOAD_SCRIPT_URL belum dikonfigurasi di Railway.", 503)
    if not secret:
        raise AppError("PHOTO_UPLOAD_SECRET belum dikonfigurasi di Railway.", 503)
    request_payload = {**payload, "secret": secret}
    relay_request = urllib.request.Request(
        script_url,
        data=json.dumps(request_payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(relay_request, timeout=45) as response:
            raw_response = response.read().decode("utf-8")
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")
        raise AppError(f"Apps Script upload relay mengembalikan HTTP {exc.code}. Detail: {detail[:500]}", 503) from exc
    except urllib.error.URLError as exc:
        raise AppError(f"Apps Script upload relay tidak dapat dihubungi. Detail: {clean(exc.reason)}", 503) from exc
    except TimeoutError as exc:
        raise AppError("Apps Script upload relay timeout. Silakan coba kembali.", 504) from exc
    try:
        result = json.loads(raw_response)
    except json.JSONDecodeError as exc:
        raise AppError(f"Respons Apps Script bukan JSON yang valid. Detail: {raw_response[:300]}", 503) from exc
    if not result.get("success"):
        raise AppError(f"Upload foto melalui Apps Script gagal. Detail: {clean(result.get('message')) or 'Tidak ada detail error.'}", 503)
    result.pop("success", None)
    return result


def get_worksheet(name):
    try:
        return get_spreadsheet().worksheet(name)
    except gspread.WorksheetNotFound as exc:
        raise AppError(f"Sheet {name} belum tersedia. Jalankan /api/setup.", 503) from exc
    except gspread.exceptions.APIError as exc:
        raise AppError(f"Sheet {name} tidak bisa dibaca. Periksa akses service account.", 503) from exc


def get_sheet_values(worksheet, sheet_name):
    try:
        return worksheet.get_all_values()
    except gspread.exceptions.APIError as exc:
        raise AppError(f"Sheet {sheet_name} tidak bisa dibaca. Periksa akses service account.", 503) from exc


def ensure_worksheet(spreadsheet, name, required_headers):
    try:
        worksheet = spreadsheet.worksheet(name)
    except gspread.WorksheetNotFound:
        worksheet = spreadsheet.add_worksheet(title=name, rows=1000, cols=max(len(required_headers), 20))
    values = get_sheet_values(worksheet, name)
    existing_headers = values[0] if values else []
    missing = [header for header in required_headers if header not in existing_headers]
    if not existing_headers:
        worksheet.update([required_headers], "A1")
    elif missing:
        start_column = len(existing_headers) + 1
        worksheet.update([missing], gspread.utils.rowcol_to_a1(1, start_column))
    worksheet.freeze(rows=1)
    worksheet.format("1:1", {"backgroundColor": {"red": 1, "green": 1, "blue": 1}, "textFormat": {"bold": True, "foregroundColor": {"red": 0, "green": 0, "blue": 0}}})
    return {"name": name, "addedHeaders": missing}


def ensure_required_headers(worksheet, sheet_name, required_headers):
    values = get_sheet_values(worksheet, sheet_name)
    existing_headers = values[0] if values else []
    missing = [header for header in required_headers if header not in existing_headers]
    if not existing_headers:
        worksheet.update([required_headers], "A1")
    elif missing:
        worksheet.update([missing], gspread.utils.rowcol_to_a1(1, len(existing_headers) + 1))
    return missing


def get_rows(worksheet, expected_headers):
    values = get_sheet_values(worksheet, worksheet.title)
    validate_headers(values, expected_headers, worksheet.title)
    return [row_to_dict(values[0], row) for row in values[1:]]


def get_rows_with_numbers(worksheet, expected_headers):
    values = get_sheet_values(worksheet, worksheet.title)
    validate_headers(values, expected_headers, worksheet.title)
    rows = []
    for row_number, values_row in enumerate(values[1:], start=2):
        row = row_to_dict(values[0], values_row)
        row["_ROW_NUMBER"] = row_number
        rows.append(row)
    return rows


def append_record(worksheet, sheet_name, required_headers, record):
    values = get_sheet_values(worksheet, sheet_name)
    validate_headers(values, required_headers, sheet_name)
    worksheet.append_row([record.get(header, "") for header in values[0]], value_input_option="USER_ENTERED")


def validate_headers(values, expected_headers, sheet_name):
    if not values:
        raise AppError(f"Header sheet {sheet_name} belum tersedia.", 503)
    missing = [header for header in expected_headers if header not in values[0]]
    if missing:
        raise AppError(f"Header sheet {sheet_name} belum lengkap: {', '.join(missing)}", 503)


def row_to_dict(headers, values):
    return dict(zip(headers, values + [""] * (len(headers) - len(values))))


def serialize_asset(row):
    return {"assetCode": clean(row["NOMOR ASSET"]), "type": clean(row["TYPE"]), "layoutNumber": clean(row["NO LAYOUT"]), "user": clean(row["USER"]), "opname": clean(row["OPNAME"]), "masterCondition": clean(row["KONDISI"]), "area": clean(row["AREA"]), "detailLocation": clean(row["LOKASI DETAIL"]), "lastCondition": clean(row["KONDISI TERAKHIR"]), "lastStatus": clean(row["STATUS TERAKHIR"]), "lastDate": clean(row["TANGGAL OPNAME TERAKHIR"]) or "-", "lastNotes": clean(row["KETERANGAN TERAKHIR"]), "lastDocumentation": clean(row["DOKUMENTASI TERAKHIR"])}


def serialize_log(row):
    return {"timestamp": clean(row["TIMESTAMP"]) or "-", "condition": clean(row["KONDISI TERAKHIR"] or row["KONDISI"]), "status": clean(row["STATUS TERAKHIR"]), "opnameDate": clean(row["TANGGAL OPNAME TERAKHIR"]) or "-", "notes": clean(row["KETERANGAN TERAKHIR"]), "documentation": clean(row["DOKUMENTASI TERAKHIR"]), "operator": clean(row["NAMA PETUGAS"]), "role": clean(row["ROLE"])}


def now_text():
    return datetime.now(ZoneInfo(TIMEZONE)).strftime("%Y-%m-%d %H:%M:%S")


def current_period_status():
    month = datetime.now(ZoneInfo(TIMEZONE)).month
    return "SUDAH OPNAME JANUARI - JUNI" if month <= 6 else "SUDAH OPNAME JULI - DESEMBER"


def normalize(value):
    return clean(value).upper()


def clean(value):
    return "" if value is None else str(value).strip()


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.getenv("PORT", "8080")))
